#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
美团黄黑编辑风 · HTML 幻灯片 → PPTX 导出
--------------------------------------------------
思路：HTML 是设计母本，用无头 Chrome 把每页渲染成高清 PNG，
再用 python-pptx 把 PNG 铺满 16:9 幻灯片（图片版，100% 保真）。

用法:
    python3 scripts/build_pptx.py                    # 渲染 slides/ 下全部页 → out/deck.pptx
    python3 scripts/build_pptx.py --slides slides --out out/deck.pptx
    python3 scripts/build_pptx.py --scale 2          # 2x 高清（默认）

依赖:
    pip install python-pptx        # 纯 Python，轻量
    Google Chrome / Chromium       # 用于渲染，脚本自动探测
"""
import argparse, glob, os, subprocess, sys, tempfile, time

SLIDE_W, SLIDE_H = 1280, 720  # 逻辑画布（16:9）

CHROME_CANDIDATES = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
    "/Applications/Brave Browser.app/Contents/MacOS/Brave Browser",
    "google-chrome", "chromium", "chromium-browser",  # Linux PATH
]


def find_chrome():
    for c in CHROME_CANDIDATES:
        if os.path.sep in c:
            if os.path.isfile(c):
                return c
        else:
            from shutil import which
            p = which(c)
            if p:
                return p
    sys.exit("✗ 未找到 Chrome/Chromium。请安装 Google Chrome 后重试。")


def render_png(chrome, html_path, png_path, scale):
    """无头 Chrome 渲染单页 HTML 为 PNG。"""
    url = "file://" + os.path.abspath(html_path)
    cmd = [
        chrome, "--headless=new", "--disable-gpu", "--hide-scrollbars",
        "--no-sandbox", "--force-color-profile=srgb",
        f"--force-device-scale-factor={scale}",
        f"--window-size={SLIDE_W},{SLIDE_H}",
        "--virtual-time-budget=5000",         # 等 Web 字体 / 图片就绪
        "--default-background-color=00000000",
        f"--screenshot={png_path}", url,
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    if not os.path.isfile(png_path):
        raise RuntimeError(f"渲染失败: {html_path}")


def build(slides_dir, out_path, scale):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        sys.exit("✗ 缺少 python-pptx。请先运行:  pip install python-pptx")

    try:
        import renumber; renumber.renumber(slides_dir)  # 导出前自动刷页码
    except Exception:
        pass
    htmls = sorted(glob.glob(os.path.join(slides_dir, "[0-9]*.html")))
    if not htmls:
        sys.exit(f"✗ {slides_dir} 下没有 .html 幻灯片")

    chrome = find_chrome()
    print(f"→ Chrome: {chrome}")
    print(f"→ 共 {len(htmls)} 页，渲染倍率 {scale}x")

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 宽屏
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmp:
        for i, html in enumerate(htmls, 1):
            png = os.path.join(tmp, f"s{i:02d}.png")
            print(f"  [{i}/{len(htmls)}] {os.path.basename(html)}")
            render_png(chrome, html, png, scale)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)

    os.makedirs(os.path.dirname(os.path.abspath(out_path)), exist_ok=True)
    prs.save(out_path)
    print(f"✓ 已导出: {out_path}")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--slides", default="slides", help="幻灯片 HTML 目录")
    ap.add_argument("--out", default="out/deck.pptx", help="输出 pptx 路径")
    ap.add_argument("--scale", type=int, default=2, help="渲染倍率（清晰度），默认 2")
    a = ap.parse_args()
    build(a.slides, a.out, a.scale)
